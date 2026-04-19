"""Generate the 5-slide pitch deck (``deck.pptx``) via python-pptx.

Slides:

1. Problem — three value streams, one battery, siloed today
2. Approach — one orchestrator, four actions, mobility risk priced
3. Math — stochastic MILP formulation, mobility-fail term highlighted
4. Results — baseline comparison + sensitivity message
5. Roadmap — pilot → OpenVPP integration → ancillary services

Pulls headline numbers from ``results/backtest_summary.csv`` at generation
time so the deck auto-updates whenever the backtest is re-run.

Run: ``python -m src.pitch.deck``.
"""

from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
RESULTS_DIR = ROOT / "results"
DECK_PATH = ROOT / "deck.pptx"

DARK = RGBColor(0x1F, 0x29, 0x37)
ACCENT = RGBColor(0xDC, 0x26, 0x26)
MUTED = RGBColor(0x6B, 0x72, 0x80)


def _clear_default_layout(slide) -> None:
    """Remove any placeholder shapes from the default slide layout."""
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)


def _text_box(slide, left, top, width, height, text, size=18, bold=False, color=DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf


def _add_bullet(tf, text: str, size: int = 16, color=DARK, bold=False) -> None:
    p = tf.add_paragraph()
    p.level = 0
    run = p.add_run()
    run.text = "• " + text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def _title(slide, text: str, subtitle: str | None = None) -> None:
    _text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.9),
              text, size=32, bold=True, color=DARK)
    if subtitle:
        _text_box(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
                  subtitle, size=16, color=MUTED)


def build_slide_problem(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _clear_default_layout(slide)
    _title(slide, "One battery, three revenue streams — siloed today.",
           "Every EV has ~100 kWh of battery, idle 18–22 h/day. Arizona grid is stressed. AI wants distributed compute.")

    tf = _text_box(slide, Inches(0.6), Inches(2.0), Inches(6), Inches(4), "", size=18)
    _add_bullet(tf, "Mobility: 25–40 kWh/day needed. Non-negotiable.", size=18, bold=True)
    _add_bullet(tf, "V2G: ~10 grid stress events/year, $20–150 revenue.", size=18)
    _add_bullet(tf, "Inference: 150–250 participation days, $400–2,000 net.", size=18)
    _add_bullet(tf, "", size=18)
    _add_bullet(tf, "No platform today co-optimizes all three.", size=18,
                bold=True, color=ACCENT)

    _text_box(slide, Inches(7.0), Inches(2.2), Inches(6), Inches(3),
              "Challenge-stated flexible capacity per vehicle: 30–50 kWh\n"
              "Arizona events: 6 summer afternoons + 4 winter mornings\n"
              "Peak LMP during events: $0.20–$0.50/kWh\n"
              "Inference revenue: $0.20–$0.60/kWh-equivalent",
              size=14, color=MUTED)


def build_slide_approach(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _clear_default_layout(slide)
    _title(slide, "One orchestrator. Four actions. Mobility risk priced, not walled.")

    tf = _text_box(slide, Inches(0.6), Inches(2.0), Inches(12), Inches(4), "", size=18)
    _add_bullet(tf,
                "Actions per vehicle per hour: charge · discharge (V2G) · run inference · idle",
                size=18, bold=True)
    _add_bullet(tf,
                "Mobility-miss penalty λ = $50/hour. The optimizer weighs missed-trip cost "
                "against grid/inference revenue instead of using a hard SoC wall.",
                size=17)
    _add_bullet(tf,
                "Rolling-horizon stochastic MILP: 24h horizon, hourly resolution, scenario-based.",
                size=17)
    _add_bullet(tf,
                "Three forecasters feed probabilistic scenarios: grid events, LMP, inference demand.",
                size=17)
    _add_bullet(tf,
                "Per-vehicle solves in parallel — architecture scales linearly to fleet size.",
                size=17)


def build_slide_math(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _clear_default_layout(slide)
    _title(slide, "The math — stochastic MILP with priced mobility risk.")

    math_text = (
        "max    E_s [ Σ_t   LMP·discharge  −  LMP·charge  +  r_inf·P_inf·infer  −  (λ/reserve)·shortfall ]\n\n"
        "s.t.   SoC[t+1,s] = SoC[t,s] + η_c·charge − (1/η_d)·discharge − P_inf·infer − trip_kwh\n"
        "       SoC ∈ [0, capacity]\n"
        "       charge, discharge, infer = 0  when vehicle not at home\n"
        "       infer ≤ demand_available\n"
        "       non-anticipativity:  charge[0,s] = charge[0,0]   ∀s   (same for discharge, infer)\n"
    )
    _text_box(slide, Inches(0.5), Inches(1.9), Inches(12.3), Inches(3.5),
              math_text, size=14, color=DARK)

    _text_box(slide, Inches(0.5), Inches(5.4), Inches(12), Inches(1.2),
              "The λ term is the differentiator. Most V2G platforms use a hard SoC floor — "
              "which either over-reserves (leaving money on the table) or under-reserves "
              "(risking mobility). Pricing the miss lets the optimizer find the economic "
              "frontier.",
              size=14, color=MUTED)


def _headline_numbers() -> dict:
    p = RESULTS_DIR / "backtest_summary.csv"
    if not p.exists():
        return {}
    s = pd.read_csv(p, index_col=0)
    out: dict = {}
    try:
        out["stochastic"] = s.loc["stochastic_coopt", "total_revenue_usd"]
        out["greedy"] = s.loc["greedy", "total_revenue_usd"]
        out["v2g_only"] = s.loc["v2g_only", "total_revenue_usd"]
        out["passive"] = s.loc["passive", "total_revenue_usd"]
        out["reliability"] = s.loc["stochastic_coopt", "mobility_reliability"] * 100
    except KeyError:
        pass
    return out


def build_slide_results(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _clear_default_layout(slide)
    _title(slide, "Results — stochastic co-opt beats every baseline on revenue.",
           "1-year backtest on a 100-vehicle synthetic Phoenix fleet.")

    nums = _headline_numbers()
    if nums:
        beat_greedy = 100 * (nums["stochastic"] - nums["greedy"]) / max(abs(nums["greedy"]), 1)
        beat_v2g = 100 * (nums["stochastic"] - nums["v2g_only"]) / max(abs(nums["v2g_only"]), 1)
        tf = _text_box(slide, Inches(0.6), Inches(1.9), Inches(6), Inches(4), "", size=18)
        _add_bullet(tf, f"Stochastic co-opt:  ${nums['stochastic']:,.0f}/veh/yr", size=20, bold=True, color=ACCENT)
        _add_bullet(tf, f"Greedy baseline:  ${nums['greedy']:,.0f}  (+{beat_greedy:.0f}% headroom)", size=18)
        _add_bullet(tf, f"V2G-only:  ${nums['v2g_only']:,.0f}", size=18)
        _add_bullet(tf, f"Passive (dumb charge):  ${nums['passive']:,.0f}", size=18)
        _add_bullet(tf, f"Mobility reliability:  {nums['reliability']:.1f}% of hours", size=18, bold=True)

    plot = RESULTS_DIR / "backtest_comparison.png"
    if plot.exists():
        slide.shapes.add_picture(str(plot), Inches(6.8), Inches(1.9), width=Inches(6.3))


def build_slide_roadmap(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _clear_default_layout(slide)
    _title(slide, "Roadmap — from simulator to pilot to platform.")

    tf = _text_box(slide, Inches(0.6), Inches(1.9), Inches(12), Inches(5), "", size=18)
    _add_bullet(tf, "Weeks 0–4: integrate with OpenVPP's DER orchestration layer — real LMP feed, driver-app opt-in",
                size=17, bold=True)
    _add_bullet(tf, "Weeks 4–12: 20-vehicle Phoenix pilot — APS TOU rate, real V2G-capable home chargers", size=17)
    _add_bullet(tf, "Quarter 2: add ancillary services (frequency regulation, voltage support) — short-horizon revenue stacks on top", size=17)
    _add_bullet(tf, "Quarter 3: EV-fleet operator SDK — dispatch API for rideshare, delivery, municipal fleets", size=17)
    _add_bullet(tf, "Defensive: sensitivity slider in the dashboard. At $0 inference revenue, we're still a smart-charging + V2G platform.",
                size=17, color=ACCENT, bold=True)


def build_deck() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    build_slide_problem(prs)
    build_slide_approach(prs)
    build_slide_math(prs)
    build_slide_results(prs)
    build_slide_roadmap(prs)
    prs.save(DECK_PATH)
    return DECK_PATH


def main() -> None:
    p = build_deck()
    print(f"deck written → {p}")


if __name__ == "__main__":
    main()
